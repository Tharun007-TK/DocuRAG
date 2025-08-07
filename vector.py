# vector.py

from langchain_ollama import OllamaEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document
import os
import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pytesseract
from PIL import Image

embeddings = OllamaEmbeddings(model="mxbai-embed-large")
db_location = "./chrome_langchain_db"
collection_name = "uploaded_docs"

vector_store = Chroma(
    collection_name=collection_name,
    persist_directory=db_location,
    embedding_function=embeddings
)

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(file):
    image = Image.open(file)
    return pytesseract.image_to_string(image)

def add_files_to_vectorstore(files):
    docs = []
    for i, file in enumerate(files):
        ext = file.name.split(".")[-1].lower()
        content = ""

        if ext == "csv":
            df = pd.read_csv(file)
            content = "\n".join(df.astype(str).apply(" ".join, axis=1))
        elif ext == "txt":
            content = file.read().decode("utf-8")
        elif ext == "pdf":
            content = extract_text_from_pdf(file)
        elif ext == "docx":
            content = extract_text_from_docx(file)
        elif ext == "pptx":
            content = extract_text_from_pptx(file)
        elif ext in ["png", "jpg", "jpeg"]:
            content = extract_text_from_image(file)
        else:
            continue

        if content.strip():
            document = Document(page_content=content, metadata={"source": file.name}, id=f"{file.name}_{i}")
            docs.append(document)

    if docs:
        vector_store.add_documents(docs)
        return True
    return False

def get_context_for_query(query, k=5):
    retriever = vector_store.as_retriever(search_kwargs={"k": k})
    results = retriever.get_relevant_documents(query)
    return "\n".join([doc.page_content for doc in results])
